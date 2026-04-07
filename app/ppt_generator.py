from __future__ import annotations

from datetime import datetime, UTC
from pathlib import Path

from pptx import Presentation


BASE_DIR = Path(__file__).resolve().parents[1]
OUTPUT_DIR = BASE_DIR / "output"


def _add_bullet_slide(prs: Presentation, title: str, lines: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()

    for idx, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        paragraph.text = line


def build_presentation(kpi: dict, forecast: dict, market: dict) -> str:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    output_path = OUTPUT_DIR / f"{kpi['brand']}_business_review.pptx"

    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = f"{kpi['brand']} Business Review"
    subtitle = title_slide.placeholders[1]
    subtitle.text = (
        f"Generated at {datetime.now(UTC).strftime('%Y-%m-%d %H:%M UTC')}\n"
        "Data source: local certified MVP data products"
    )

    _add_bullet_slide(
        prs,
        "KPI Snapshot",
        [
            f"Reporting period: {kpi['period']}",
            f"Sales total: {kpi['sales_total']:,.0f}",
            f"Sales growth: {kpi['sales_growth_pct'] * 100:.1f}%",
            f"Conversion rate: {kpi['conversion_rate'] * 100:.1f}%",
            f"Market share: {kpi['market_share'] * 100:.1f}%",
        ],
    )

    _add_bullet_slide(
        prs,
        "Forecast Slide",
        [
            f"Product: {forecast['product']}",
            f"Forecast EoM: {forecast['forecast_eom']:,.0f}",
            f"Forecast tercial / next quarter total: {forecast['forecast_next_quarter_total']:,.0f}",
            *[
                (
                    f"{row['month']}: {row['forecast_sales']:,.0f} "
                    f"(CI {row['lower_ci']:,.0f} - {row['upper_ci']:,.0f})"
                )
                for row in forecast["forecast_months"]
            ],
        ],
    )

    _add_bullet_slide(
        prs,
        "Market Development",
        [
            f"Latest market size: {market['latest_market_size']:,.0f}",
            f"Market size growth vs. first month: {market['market_size_growth_pct'] * 100:.1f}%",
            f"Latest market growth input: {market['latest_market_growth_pct'] * 100:.1f}%",
            "Narrative: Market continues to expand and supports sustained sales growth assumptions.",
        ],
    )

    prs.save(output_path)
    return str(output_path)
